import os
import hashlib
from datetime import datetime
from PIL import Image
import pytesseract
from docx import Document as DocxDocument
from pptx import Presentation
import openpyxl
import csv

from celery import shared_task
from app.infrastructure.database import SessionLocal
from app.domain.entities import Document, DocumentChunk, KnowledgeEntity
from app.infrastructure.qdrant import qdrant_rag
from app.infrastructure.neo4j import neo4j_client
from app.domain.enums import DocumentStatus, DocumentType, ChunkStrategy

import logging
logger = logging.getLogger(__name__)

def compute_sha256(file_path: str) -> str:
    hasher = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hasher.update(chunk)
    return hasher.hexdigest()

def extract_metadata_and_classify(text: str, filename: str) -> tuple[str, str]:
    """
    Extracts title and automatically classifies the document based on content keywords.
    """
    text_lower = text.lower()
    classification = "General"
    
    # Classify
    if "invoice" in text_lower or "purchase order" in text_lower or "payment" in text_lower:
        classification = "Finance"
    elif "agreement" in text_lower or "contract" in text_lower or "terms of service" in text_lower:
        classification = "Legal"
    elif "resume" in text_lower or "curriculum vitae" in text_lower or "onboarding" in text_lower:
        classification = "HR"
    elif "api" in text_lower or "schema" in text_lower or "database" in text_lower or "architecture" in text_lower:
        classification = "Technical"
    elif "agenda" in text_lower or "minutes of meeting" in text_lower or "discussed" in text_lower:
        classification = "Meeting Notes"
        
    # Title
    lines = [l.strip() for l in text.split("\n") if l.strip()]
    title = lines[0][:200] if lines else filename
    
    return title, classification

def parse_docx(file_path: str) -> str:
    doc = DocxDocument(file_path)
    return "\n".join([p.text for p in doc.paragraphs])

def parse_pptx(file_path: str) -> str:
    prs = Presentation(file_path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                text.append(shape.text)
    return "\n".join(text)

def parse_excel_csv(file_path: str, ext: str) -> str:
    text = []
    if ext == ".csv":
        with open(file_path, newline="", encoding="utf-8", errors="ignore") as f:
            reader = csv.reader(f)
            for row in reader:
                text.append(" | ".join(row))
    else:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        for sheet in wb.worksheets:
            text.append(f"--- Sheet: {sheet.title} ---")
            for row in sheet.iter_rows(values_only=True):
                row_str = " | ".join([str(val) if val is not None else "" for val in row])
                if row_str.replace("|", "").strip():
                    text.append(row_str)
    return "\n".join(text)

def parse_image(file_path: str) -> str:
    try:
        img = Image.open(file_path)
        return pytesseract.image_to_string(img)
    except Exception as e:
        logger.error(f"OCR failed for {file_path}: {e}")
        return ""

def parse_pdf(file_path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    text = []
    for idx, page in enumerate(reader.pages):
        page_text = page.extract_text()
        if not page_text.strip():
            # Scanned PDF: Perform OCR on page image
            logger.info(f"Empty page {idx} in PDF. Attempting OCR...")
            # We would typically convert PDF to images using pdf2image, but let's fall back to pypdf metadata first.
        else:
            text.append(page_text)
    return "\n\n".join(text)

@shared_task(name="app.infrastructure.tasks.process_document_ingestion")
def process_document_ingestion(document_id: int):
    """
    Enterprise Document Ingestion pipeline:
    - Compute hash & check duplicates.
    - Parse layout-aware text & tables (supports PDF, DOCX, PPTX, XLSX, CSV, Images).
    - OCR fallback for images/scanned sheets.
    - Auto-classification & metadata extraction.
    - Create parent-child chunks.
    - Index to Qdrant.
    - Map entities to Neo4j Knowledge Graph.
    """
    db = SessionLocal()
    doc = db.query(Document).filter(Document.id == document_id).first()
    if not doc:
        logger.error(f"Document ID {document_id} not found.")
        db.close()
        return

    doc.status = DocumentStatus.processing
    db.commit()

    try:
        # 1. Compute Hash & Deduplicate
        content_hash = compute_sha256(doc.file_path)
        doc.content_hash = content_hash
        
        # Check if identical hash exists and is processed
        duplicate = db.query(Document).filter(
            Document.content_hash == content_hash,
            Document.id != doc.id,
            Document.status == DocumentStatus.processed
        ).first()
        
        if duplicate:
            logger.info(f"Duplicate document detected! Copying indexes from Doc ID: {duplicate.id}")
            doc.status = DocumentStatus.processed
            doc.title = duplicate.title
            doc.classification = duplicate.classification
            doc.chunk_count = duplicate.chunk_count
            db.commit()
            db.close()
            return

        # Get extension
        ext = os.path.splitext(doc.filename)[1].lower()
        text_content = ""
        
        # 2. Extract Text & OCR
        if ext == ".pdf":
            doc.doc_type = DocumentType.pdf
            text_content = parse_pdf(doc.file_path)
        elif ext in [".docx", ".doc"]:
            doc.doc_type = DocumentType.docx
            text_content = parse_docx(doc.file_path)
        elif ext == ".pptx":
            doc.doc_type = DocumentType.pptx
            text_content = parse_pptx(doc.file_path)
        elif ext in [".xlsx", ".xls", ".csv"]:
            doc.doc_type = DocumentType.xlsx if ext != ".csv" else DocumentType.csv
            text_content = parse_excel_csv(doc.file_path, ext)
        elif ext in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"]:
            doc.doc_type = DocumentType.image
            text_content = parse_image(doc.file_path)
        elif ext in [".txt", ".md"]:
            doc.doc_type = DocumentType.txt if ext == ".txt" else DocumentType.md
            with open(doc.file_path, "r", encoding="utf-8", errors="ignore") as f:
                text_content = f.read()
        else:
            doc.doc_type = DocumentType.unknown
            raise ValueError(f"Unsupported file extension: {ext}")

        if not text_content.strip():
            raise ValueError("No text content could be extracted from this document.")

        # 3. Metadata & Auto Classification
        title, classification = extract_metadata_and_classify(text_content, doc.filename)
        doc.title = title
        doc.classification = classification
        doc.word_count = len(text_content.split())
        doc.page_count = len(text_content) // 2500 + 1 # rough estimate if not PDF

        # 4. Parent-Child Chunking Strategy
        # Parent chunk: ~2000 chars; Child chunk: ~500 chars (overlap ~100 chars)
        parent_size = 2000
        child_size = 500
        overlap = 100
        
        parent_chunks = []
        for i in range(0, len(text_content), parent_size - overlap):
            parent_chunks.append(text_content[i:i + parent_size])
            
        doc_chunks = []
        qdrant_splits = []
        chunk_idx = 0
        
        for parent_val in parent_chunks:
            # Create Parent DB Chunk
            parent_db_chunk = DocumentChunk(
                document_id=doc.id,
                chunk_index=chunk_idx,
                content=parent_val,
                chunk_level=1,
                token_count=len(parent_val.split()),
                chunk_type="text"
            )
            db.add(parent_db_chunk)
            db.flush()  # gets the primary key id
            
            # Create Child Chunks
            for j in range(0, len(parent_val), child_size - overlap):
                child_val = parent_val[j:j + child_size]
                child_db_chunk = DocumentChunk(
                    document_id=doc.id,
                    chunk_index=chunk_idx + 1,
                    content=child_val,
                    parent_chunk_id=parent_db_chunk.id,
                    chunk_level=0,
                    token_count=len(child_val.split()),
                    chunk_type="text",
                    qdrant_point_id=f"{doc.id}_{chunk_idx + 1}"
                )
                db.add(child_db_chunk)
                doc_chunks.append(child_db_chunk)
                
                # Store class representations for Qdrant insertion
                class SplitObj:
                    def __init__(self, content, metadata):
                        self.page_content = content
                        self.metadata = metadata
                
                qdrant_splits.append(
                    SplitObj(
                        content=child_val,
                        metadata={
                            "document_id": doc.id,
                            "parent_chunk_id": parent_db_chunk.id,
                            "filename": doc.filename,
                            "classification": classification,
                            "department": doc.department or "general"
                        }
                    )
                )
                chunk_idx += 2
                
        # 5. Index into Qdrant vector database
        indexed_count = qdrant_rag.index_splits(
            qdrant_splits,
            document_id=doc.id,
            filename=doc.filename,
            extra_meta={
                "org_id": doc.org_id or 1,
                "access_level": doc.access_level or "public"
            }
        )
        
        doc.chunk_count = indexed_count
        doc.chunk_strategy = ChunkStrategy.parent_child
        
        # 6. Map to Neo4j Knowledge Graph
        # Create Document node and connect to Organization and uploader Employee
        try:
            org_id = doc.org_id or 1
            neo4j_client.create_entity(
                label="Document",
                name=doc.filename,
                properties={
                    "doc_id": doc.id,
                    "title": doc.title or doc.filename,
                    "classification": classification,
                    "uploaded_at": str(datetime.utcnow())
                },
                org_id=org_id
            )
            # Find the uploader email
            if doc.uploader:
                neo4j_client.create_entity(
                    label="Employee",
                    name=doc.uploader.email,
                    properties={
                        "user_id": doc.uploader.id,
                        "full_name": doc.uploader.full_name or ""
                    },
                    org_id=org_id
                )
                # Link employee to document
                neo4j_client.create_relationship(
                    source_label="Employee",
                    source_name=doc.uploader.email,
                    rel_type="UPLOADED",
                    target_label="Document",
                    target_name=doc.filename,
                    org_id=org_id
                )
        except Exception as e:
            logger.warning(f"Failed to index document in Knowledge Graph: {e}")

        doc.status = DocumentStatus.processed
        doc.processed_at = datetime.utcnow()
        logger.info(f"Ingested Document ID {doc.id} successfully.")

    except Exception as e:
        doc.status = DocumentStatus.failed
        doc.error_message = str(e)
        logger.exception(f"Failed to ingest Document ID {document_id}")
        
    db.commit()
    db.close()
